"""Client-side report rendering (PPTX + PDF).

Per section 9 of the architecture:
  1. Client builds a findings payload (questions, answer text, column names,
     code snippets, has_chart/has_table) — no data values.
  2. Client sends the findings payload to the brain.
  3. Brain runs the LLM to produce the narrative content (report_structure)
     and returns it as structured content.
  4. Client renders charts locally (kaleido) and merges the returned content
     into ITS OWN template via python-pptx, producing the final file locally.

The template / brand colors / fonts NEVER leave this server.
"""
from __future__ import annotations

import base64
import io
import json
import re
import secrets
from datetime import datetime
from pathlib import Path
from typing import Any

from fastapi import APIRouter, Request
from fastapi.responses import JSONResponse, Response

import brain_client
import local_store
from brain_client import TenantRevokedError, BrainError
from logger_utils import log_with_sid

# Module-level `qn` (OOXML qualified-name helper). Several helpers
# (_list_native_bullet / _suppress_paragraph_bullet) resolve bullet markers via
# `qn(...)`; without this they raised AttributeError that their try/except
# swallowed, so native markers were never detected and our explicit enumerator
# stacked on top (PROBLEM 1). _clone_slide also re-imports it locally — harmless.
from pptx.oxml.ns import qn


router = APIRouter(prefix="/api/chat", tags=["client-report"])


def _is_plotly_html(image_b64) -> bool:
    if not image_b64:
        return False
    s = str(image_b64).strip()
    if s.startswith("<") and "plotly" in s[:1000].lower():
        return True
    return s.startswith("<!") or s.startswith("<div")


def _plotly_html_to_png(html_str: str, sid: str) -> bytes | None:
    """Mirror chat.py.plotly_html_to_png — local, uses kaleido. No network."""
    try:
        import plotly.graph_objects as _go

        header = re.search(r'Plotly\.(?:newPlot|react)\s*\(\s*["\']([^"\']+)["\']\s*,\s*', html_str)
        if not header:
            return None
        # Reuse the simple balanced-bracket scan from chat.py
        def _find_balanced_end(s, start):
            if start >= len(s):
                return -1
            open_ch = s[start]
            close_ch = ']' if open_ch == '[' else '}'
            depth = 0
            i = start
            in_string = False
            string_quote = ''
            while i < len(s):
                ch = s[i]
                if in_string:
                    if ch == '\\':
                        i += 2
                        continue
                    if ch == string_quote:
                        in_string = False
                    i += 1
                    continue
                if ch in ('"', "'"):
                    in_string = True
                    string_quote = ch
                    i += 1
                    continue
                if ch == open_ch:
                    depth += 1
                elif ch == close_ch:
                    depth -= 1
                    if depth == 0:
                        return i + 1
                i += 1
            return -1

        data_start = header.end()
        if data_start >= len(html_str) or html_str[data_start] != '[':
            return None
        data_end = _find_balanced_end(html_str, data_start)
        if data_end < 0:
            return None
        data_json = json.loads(html_str[data_start:data_end])
        i = data_end
        while i < len(html_str) and html_str[i] in ' \t\n\r,':
            i += 1
        if i >= len(html_str) or html_str[i] != '{':
            return None
        layout_end = _find_balanced_end(html_str, i)
        layout_json = json.loads(html_str[i:layout_end])
        fig = _go.Figure(data=data_json, layout=layout_json)
        return fig.to_image(format="png", width=1200, height=800, scale=2)
    except Exception as e:
        log_with_sid(sid, "warning", f"PLOTLY_TO_PNG_FAILED: {e}")
        return None


def _resolve_chart_image(image_b64, sid):
    if not image_b64:
        return None, False
    if _is_plotly_html(image_b64):
        png = _plotly_html_to_png(image_b64, sid)
        if png:
            return base64.b64encode(png).decode("ascii"), True
        return None, True
    return image_b64, True


def _build_qa_pairs(history: list, sid: str):
    """Same pairing logic as B2C `_build_report_data`."""
    qa_pairs = []
    i = 0
    while i < len(history):
        msg = history[i]
        if msg.get("role") == "human":
            if i + 1 < len(history) and history[i + 1].get("role") == "ai":
                ai_msg = history[i + 1]
                images_array = ai_msg.get("images") or []
                if images_array:
                    for img_entry in images_array:
                        resolved_img, is_chart = _resolve_chart_image(img_entry.get("image_base64", ""), sid)
                        qa_pairs.append({
                            "human": msg,
                            "ai": {"content": img_entry.get("answer", ""), "image_base64": resolved_img, "_resolved_png": resolved_img},
                            "ai_index": i + 1,
                            "has_chart": is_chart,
                            "has_table": False,
                        })
                else:
                    resolved_img, is_chart = _resolve_chart_image(ai_msg.get("image_base64", ""), sid)
                    has_table = bool(ai_msg.get("table"))
                    has_text = bool((ai_msg.get("content") or "").strip())
                    ai_copy = dict(ai_msg)
                    if is_chart:
                        ai_copy["_resolved_png"] = resolved_img
                    if is_chart or has_table or has_text:
                        qa_pairs.append({
                            "human": msg, "ai": ai_copy, "ai_index": i + 1,
                            "has_chart": is_chart, "has_table": has_table,
                        })
                i += 2
                continue
        i += 1
    return qa_pairs


def _sanitize_filename(name: str) -> str:
    name = re.sub(r"[^a-zA-Z0-9_-]", "_", (name or "report").strip())
    return name[:30] or "report"


def _build_findings_for_brain(qa_pairs: list) -> list[dict]:
    """The findings shape the brain expects on /v1/report — same shape the
    B2C `_generate_report_structure` built internally."""
    out = []
    for idx, pair in enumerate(qa_pairs):
        out.append({
            "index": idx,
            "question": (pair["human"].get("content") or "")[:500],
            "answer_text": (pair["ai"].get("content") or "")[:500],
            "has_chart": bool(pair["has_chart"]),
            "has_table": bool(pair["has_table"]),
            "table_columns": ((pair["ai"].get("table") or {}).get("columns") or [])[:10],
            "code_snippet": (pair["ai"].get("code") or "")[:300],
        })
    return out


def _slot_budgets_for(sid: str) -> dict:
    """Per-slot character budgets from the tenant's design spec so the brain
    generates titles/bodies that FIT their slots (Phase 1/3). Prefers an explicit
    `max_chars` constraint from the analyzer's layout_plan; else estimates from
    slot geometry with the same char-box model the renderer uses. Generic — keys
    off the content/cover title + body regions, never template-specific names.
    Returns {} when no usable plan exists (the brain then uses sane defaults)."""
    def _cap(region, size):
        if not isinstance(region, dict):
            return None
        mc = region.get("max_chars")
        if isinstance(mc, (int, float)) and mc > 0:
            return int(mc)
        try:
            w = float(region.get("w") or 0)
            h = float(region.get("h") or 0)
        except Exception:
            return None
        if w <= 0 or h <= 0:
            return None
        size_in = max(float(size), 1.0) / 72.0
        cpl = max(1, int(w / (0.5 * size_in)))
        lines = max(1, int(h / (1.2 * size_in)))
        return cpl * lines
    try:
        import pptx_template_cache
        bundle = pptx_template_cache.get_template_and_spec(sid)
        slides = ((bundle or {}).get("layout_plan") or {}).get("slides") or {}
        content = slides.get("content") or {}
        cover = slides.get("cover") or {}
        budgets: dict = {}
        c_title = content.get("title") or {}
        ptitle = _cap(c_title, float(c_title.get("size_pt") or 26))
        body = _cap(content.get("body"), 9.0)
        cv_title = cover.get("title") or {}
        rtitle = _cap(cv_title, float(cv_title.get("size_pt") or 32))
        if ptitle:
            budgets["page_title"] = min(ptitle, 56)
        if body:
            budgets["body"] = min(body, 700)
        if rtitle:
            budgets["report_title"] = min(rtitle, 80)
        return budgets
    except Exception as e:
        log_with_sid(sid, "info", "SLOT_BUDGETS_DEFAULT", reason=type(e).__name__)
        return {}


def _build_payload(request: Request, chat_id: str, conv_id: str):
    """Shared validation + narrative build. Returns (qa_pairs, report_structure) or JSONResponse."""
    email = request.session.get("email")
    if not email:
        return JSONResponse({"error": "Not authenticated"}, status_code=401)
    if not local_store.chat_exists(chat_id):
        return JSONResponse({"error": "Chat not found"}, status_code=404)
    # Authorize owner OR a shared recipient — mirror routes/chat.py `_require_chat`
    # so report download (PPTX/PDF) is allowed for everyone who can view the chat
    # (chat-level + conversation-level sharing both populate
    # meta.json["sharing"]["shared_with"]). A non-owner may additionally export
    # ONLY a conversation that exists in their OWN index — never the owner's
    # other conversations in the same shared chat.
    if local_store.get_chat_meta_owner(chat_id) != email:
        allowed = False
        try:
            store = local_store.ChatDataStore(chat_id)
            sharing = (store.read_meta().get("sharing") or {}).get("shared_with") or []
            if email in [s.lower() for s in sharing]:
                allowed = True
        except Exception:
            allowed = False
        if not allowed:
            return JSONResponse({"error": "Access denied"}, status_code=403)
        if not local_store.user_owns_conversation(email, chat_id, conv_id):
            return JSONResponse({"error": "Access denied"}, status_code=403)

    sid = secrets.token_hex(8)
    store = local_store.ChatDataStore(chat_id)
    history = store.get_history(conv_id)
    qa_pairs = _build_qa_pairs(history, sid)
    if not qa_pairs:
        return JSONResponse({"error": "No reportable content"}, status_code=400)

    findings = _build_findings_for_brain(qa_pairs)
    try:
        rsp = brain_client.report(sid=sid, qa_pairs=findings, user_email=email,
                                  slot_budgets=_slot_budgets_for(sid))
    except TenantRevokedError:
        return JSONResponse({"error": "Service unavailable"}, status_code=503)
    except BrainError as e:
        log_with_sid(sid, "error", f"REPORT_BRAIN_ERROR: {e}")
        return JSONResponse({"error": "Could not generate narrative"}, status_code=502)
    return qa_pairs, rsp.get("report_structure") or {}, sid


@router.post("/{chat_id}/conversation/{conv_id}/download_report")
async def download_pdf(request: Request, chat_id: str, conv_id: str):
    """PDF export — uses reportlab to render the brain's narrative output
    into a branded document. Same color palette / fonts as the B2C app."""
    res = _build_payload(request, chat_id, conv_id)
    if isinstance(res, JSONResponse):
        return res
    qa_pairs, report_structure, sid = res

    pdf_bytes = _render_pdf(qa_pairs, report_structure, sid)
    filename = _sanitize_filename(report_structure.get("filename") or "analysis") + ".pdf"
    try:
        import brain_client
        brain_client.post_activity("report_exported",
                                   request.session.get("email", ""),
                                   {"format": "pdf", "chat_id": chat_id, "conv_id": conv_id})
    except Exception:
        pass
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/{chat_id}/conversation/{conv_id}/download_pptx")
async def download_pptx(request: Request, chat_id: str, conv_id: str):
    res = _build_payload(request, chat_id, conv_id)
    if isinstance(res, JSONResponse):
        return res
    qa_pairs, report_structure, sid = res
    pptx_bytes = _render_pptx(qa_pairs, report_structure, sid)
    filename = _sanitize_filename(report_structure.get("filename") or "analysis") + ".pptx"
    try:
        import brain_client
        brain_client.post_activity("report_exported",
                                   request.session.get("email", ""),
                                   {"format": "pptx", "chat_id": chat_id, "conv_id": conv_id})
    except Exception:
        pass
    return Response(
        content=pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


def _render_pptx(qa_pairs: list, report_structure: dict, sid: str) -> bytes:
    """PPTX rendering entry point — most-faithful → least, never crash (Art. IV).

    Fallback chain when a tenant has uploaded a branded `.pptx`:

      1. TEMPLATED CLONE (most faithful) — when the brain's v2 build-plan deck
         is usable. Deep-clone the tenant's DESIGNED slides (so backgrounds,
         header/footer art, dividers, logos and theme all carry through), drop
         the template author's own sample shapes, inject our title / narrative /
         agenda / takeaways into the template's designated shapes (shrink-to-fit),
         and place the chart in the analyzer-chosen region. See
         `_render_pptx_templated_clone`.
      2. DESIGN-SPEC NATIVE — when only a usable version-3 `layout_plan.json`
         is available. Fresh editable deck on the template base with native
         textboxes/pictures at the plan's regions. See `_render_pptx_native`.
      3. BUILT-IN — no template / nothing usable. Navy #001E44 / blue #1276C2
         PowerDataChat-branded deck. Unchanged. See `_render_pptx_builtin`.

    Templated decks (1 and 2) intentionally omit the PowerDataChat logo; only
    the tenant's branding shows.
    """
    try:
        import pptx_template_cache
        bundle = pptx_template_cache.get_template_and_spec(sid)
    except Exception as e:
        log_with_sid(sid, "warning", f"PPTX_TEMPLATE_LOOKUP_FAIL: {e}")
        bundle = {"has_template": False, "spec": None, "template_path": None,
                  "layout_plan": None}

    spec = bundle.get("spec")
    layout_plan = bundle.get("layout_plan")
    has_tpl = bool(bundle.get("has_template") and bundle.get("template_path"))
    plan_usable = _layout_plan_usable(layout_plan)
    deck_usable = _spec_deck_usable(spec)

    # 1) Templated clone — most faithful.
    if has_tpl and deck_usable:
        deck = (spec or {}).get("deck") or {}
        log_with_sid(sid, "info", "PPTX_RENDER_TEMPLATED", mode="clone",
                     spec_version=(spec or {}).get("version"),
                     cover=deck.get("cover_slide_index"),
                     agenda=deck.get("agenda_slide_index"),
                     content=deck.get("content_slide_index"))
        try:
            return _render_pptx_templated_clone(qa_pairs, report_structure, sid,
                                                bundle["template_path"], spec, layout_plan)
        except Exception as e:
            log_with_sid(sid, "warning",
                         f"PPTX_CLONE_RENDER_FAIL — falling back to native: {e}")
            log_with_sid(sid, "info", "PPTX_TPL_FALLBACK",
                         reason=f"clone_exception:{type(e).__name__}")

    # 2) Design-spec native render (fallback).
    if has_tpl and plan_usable:
        log_with_sid(sid, "info", "PPTX_RENDER_TEMPLATED",
                     mode="design_spec_native",
                     plan_version=(layout_plan or {}).get("version"),
                     n_branding=len((layout_plan or {}).get("branding") or []))
        try:
            return _render_pptx_native(qa_pairs, report_structure, sid,
                                       bundle["template_path"], layout_plan)
        except Exception as e:
            log_with_sid(sid, "warning",
                         f"PPTX_TEMPLATED_RENDER_FAIL — falling back to built-in: {e}")
            log_with_sid(sid, "info", "PPTX_TPL_FALLBACK",
                         reason=f"render_exception:{type(e).__name__}")
    else:
        log_with_sid(sid, "info", "PPTX_RENDER_BUILTIN",
                     has_template=bundle.get("has_template"),
                     has_plan=isinstance(layout_plan, dict),
                     plan_usable=plan_usable, deck_usable=deck_usable,
                     plan_version=(layout_plan or {}).get("version"))
        if bundle.get("has_template") and not plan_usable and not deck_usable:
            log_with_sid(sid, "info", "PPTX_TPL_FALLBACK",
                         reason="layout_plan_not_usable")
    return _render_pptx_builtin(qa_pairs, report_structure, sid)


def _spec_deck_usable(spec) -> bool:
    """Mirror of `brain.pptx_template_analyzer.spec_is_usable` — a usable v2
    build plan carries integer cover + content indexes, those slide keys, and a
    content chart_region. Gates the clone renderer; on False we fall through to
    the native / built-in renderers."""
    if not isinstance(spec, dict):
        return False
    if int(spec.get("version") or 0) != 2:
        return False
    deck = spec.get("deck") or {}
    cover = deck.get("cover_slide_index")
    content = deck.get("content_slide_index")
    if not isinstance(cover, int) or not isinstance(content, int):
        return False
    slides = spec.get("slides") or {}
    if str(cover) not in slides or str(content) not in slides:
        return False
    if not isinstance(slides[str(content)].get("chart_region"), dict):
        return False
    return True


def _render_pptx_builtin(qa_pairs: list, report_structure: dict, sid: str) -> bytes:
    """Built-in (no-template) PPTX renderer — navy / blue PowerDataChat brand.

    This is the unchanged general renderer that runs whenever a tenant has not
    uploaded a branded `.pptx`. The PowerDataChat logo is kept here; templated
    decks omit it.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(0x00, 0x1E, 0x44)

    def _add_title(slide, text, size=24):
        tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.5), Inches(0.8))
        tf = tx.text_frame
        tf.text = text or ""
        p = tf.paragraphs[0]
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = "DejaVu Sans"
        return tx

    def _add_accent(slide, y=1.0):
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(y), Inches(2), Inches(0.05))

    def _add_text(slide, text, top=1.3, size=14, color=None):
        tx = slide.shapes.add_textbox(Inches(0.4), Inches(top), Inches(12.5), Inches(5))
        tf = tx.text_frame
        tf.word_wrap = True
        tf.text = text or ""
        for p in tf.paragraphs:
            p.font.size = Pt(size)
            p.font.name = "DejaVu Sans"
            if color is not None:
                p.font.color.rgb = color
        return tx

    blank_layout = prs.slide_layouts[6]

    # Cover slide
    cover = prs.slides.add_slide(blank_layout)
    _add_title(cover, report_structure.get("report_title") or "Data Analysis Report", size=40)
    _add_accent(cover, y=1.4)
    _add_text(cover, report_structure.get("executive_summary") or "", top=2.0, size=14)
    _add_text(
        cover,
        f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}",
        top=6.7, size=11, color=RGBColor(0x6b, 0x72, 0x80),
    )

    # Per-finding slides
    findings = report_structure.get("findings") or []
    for i, pair in enumerate(qa_pairs):
        slide = prs.slides.add_slide(blank_layout)
        narrative = findings[i].get("narrative") if i < len(findings) else (pair["ai"].get("content") or "")
        question = (pair["human"].get("content") or "") if isinstance(pair.get("human"), dict) else ""
        page_title = _finding_title(findings[i].get("page_title") if i < len(findings) else "", question)
        _add_title(slide, _strip_emphasis(page_title))
        _add_accent(slide, y=1.05)

        # Body — narrative (strip **emphasis** markers; built-in deck is plain)
        _add_text(slide, _strip_emphasis(narrative), top=1.3, size=14)

        # Chart
        chart_b64 = pair["ai"].get("_resolved_png") or pair["ai"].get("image_base64")
        if chart_b64 and not _is_plotly_html(chart_b64):
            try:
                img_bytes = base64.b64decode(chart_b64)
                slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.4), Inches(3.0), width=Inches(12.5), height=Inches(4.2))
            except Exception as e:
                log_with_sid(sid, "warning", f"PPTX_IMG_FAIL: {e}")

        # Table
        table = pair["ai"].get("table")
        if table and not chart_b64:
            cols = table.get("columns") or []
            rows = (table.get("rows") or [])[:15]
            if cols and rows:
                tbl = slide.shapes.add_table(len(rows) + 1, len(cols), Inches(0.4), Inches(3.0), Inches(12.5), Inches(4.2)).table
                for j, c in enumerate(cols):
                    tbl.cell(0, j).text = str(c)
                for r_idx, row in enumerate(rows):
                    for j, c in enumerate(cols):
                        tbl.cell(r_idx + 1, j).text = str(row.get(c, ""))[:80]

    # Key takeaways
    takeaways = report_structure.get("key_takeaways") or []
    if len(qa_pairs) >= 3 and takeaways:
        slide = prs.slides.add_slide(blank_layout)
        _add_title(slide, "Key Takeaways")
        _add_accent(slide, y=1.05)
        bullet_text = "\n\n".join(f"• {_strip_emphasis(t)}" for t in takeaways)
        _add_text(slide, bullet_text, top=1.5, size=15)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _hex_to_rgb(hex_str: str, fallback=(55, 65, 81)):
    """Convert "RRGGBB" (no '#') to (r, g, b). Falls back to a neutral gray."""
    from pptx.dml.color import RGBColor
    try:
        s = (hex_str or "").strip().lstrip("#")
        if len(s) == 6:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        pass
    return RGBColor(*fallback)


# ---------------------------------------------------------------------------
# Emphasis + real-title helpers (iteration 5)
# The brain marks 1-3 key terms per line with **double asterisks**; we render
# those spans in bold + the theme accent and STRIP the markers so literal
# asterisks never appear. We also derive a real 3-4 word title from the
# question whenever the brain falls back to a "Finding N" placeholder.
# ---------------------------------------------------------------------------
_EMPH_RE = re.compile(r"\*\*(.+?)\*\*")
_FINDING_N_RE = re.compile(r"^\s*finding\s+\d+\s*$", re.IGNORECASE)
# Lone *italic* markers the model sometimes adds in ADDITION to **bold** (e.g. a
# lead verb "*Plots* ..."). We do NOT bold these (the accent is reserved for the
# **double-asterisk** key terms); we only strip the stray markers so a literal
# '*' never reaches the slide. The look-arounds avoid touching '**' pairs.
_SINGLE_EMPH_RE = re.compile(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)")


def _strip_emphasis(text: str) -> str:
    """Remove emphasis markers, keeping inner text. Handles BOTH **bold** and
    lone *italic* so neither ever renders as a literal asterisk. Never None."""
    t = _EMPH_RE.sub(r"\1", text or "")
    t = _SINGLE_EMPH_RE.sub(r"\1", t)
    return t


def _parse_emphasis(line: str):
    """Split `line` into [(text, is_bold), ...] on **...** spans.

    Caps emphasis at 3 bold spans (extra markers are treated as plain text),
    drops empty segments, and returns [(line, False)] when there are no markers.
    Callers wrap this in try/except so a malformed line falls back to plain."""
    s = line or ""
    if "**" not in s:
        # No bold spans — still strip any stray single-asterisk markers.
        return [(_strip_emphasis(s), False)]
    segs = []
    pos = 0
    bold_count = 0
    for m in _EMPH_RE.finditer(s):
        if bold_count >= 3:
            break
        if m.start() > pos:
            pre = _strip_emphasis(s[pos:m.start()])   # plain piece: drop lone '*'
            if pre:
                segs.append((pre, False))
        inner = m.group(1)
        if inner:
            segs.append((inner, True))
            bold_count += 1
        pos = m.end()
    if pos < len(s):
        tail = _strip_emphasis(s[pos:])  # strip any leftover markers beyond cap
        if tail:
            segs.append((tail, False))
    if not segs:
        return [(_strip_emphasis(s), False)]
    return segs


def _short_title_from_question(q: str, max_words: int = 4) -> str:
    """Derive a 3-4 word Title-Case slide title from a question (local port of
    the brain helper, Article XI — small duplication is fine). Falls back to
    'Data Insights' when nothing meaningful remains; caps at 60 chars."""
    stop = {
        "the", "a", "an", "of", "in", "on", "for", "to", "and", "or", "by",
        "with", "show", "me", "what", "is", "are", "how", "which", "per",
        "vs", "versus", "do", "does", "my", "our", "this", "that", "these",
        "those", "please", "give", "display", "plot", "chart", "graph", "top",
        "between", "from",
    }
    try:
        tokens = re.findall(r"[A-Za-z][A-Za-z0-9']*", q or "")
        kept = []
        for tok in tokens:
            if tok.lower() in stop:
                continue
            kept.append(tok if (tok.isupper() and len(tok) >= 2) else tok.capitalize())
            if len(kept) >= max_words:
                break
        title = " ".join(kept).strip()
        if not title:
            return "Data Insights"
        return title[:60].rstrip()
    except Exception:
        return "Data Insights"


def _finding_title(title: str, question: str) -> str:
    """Return a real slide title, HARD-CAPPED to 4 words (iteration 6) so it fits
    the header on one line and is never ellipsis-truncated. If `title` is empty
    or a 'Finding N' placeholder, derive a <=4-word Title-Case title from
    `question` (the brain already caps to 4 words; this is a defensive backstop
    covering the passthrough case)."""
    t = (title or "").strip()
    if not t or _FINDING_N_RE.match(t):
        return _short_title_from_question(question or "")
    words = t.split()
    if len(words) > 4:
        t = " ".join(words[:4]).rstrip(" .,:;-—–") or t
    return t


def _write_runs(p, line: str, *, size_pt, font_name, color_rgb, bold,
                emphasis_color=None) -> None:
    """Write `line` into paragraph `p` as one or more runs.

    When `emphasis_color` is given, **...** spans become bold runs in that color
    (markers stripped); other segments use the base bold/color. When it is None,
    a single run is written with the markers stripped. Always sets the font name
    + size; color assignment is wrapped (some run states reject it)."""
    from pptx.util import Pt
    if emphasis_color is not None:
        try:
            segments = _parse_emphasis(line)
        except Exception:
            segments = [(_strip_emphasis(line), False)]
    else:
        segments = [(_strip_emphasis(line), False)]
    for text, is_bold in segments:
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        if is_bold:
            run.font.bold = True
            try:
                run.font.color.rgb = emphasis_color
            except Exception:
                pass
        else:
            run.font.bold = bold
            try:
                run.font.color.rgb = color_rgb
            except Exception:
                pass


# Leading enumerator / bullet at the start of a list line: "1. ", "1) ",
# "• ", "- ", "* ", "· ", etc. Used to guarantee EXACTLY ONE marker per
# agenda/takeaways line (template native bullet OR explicit, never both).
_LEADING_LIST_MARKER_RE = re.compile(
    r"^\s*(?:\d+[.)]|[•‣◦⁃∙·\-\*])\s+"
)


def _strip_leading_list_marker(text: str) -> str:
    """Remove a single leading enumerator/bullet ('1. ', '1) ', '• ', '- ',
    '* ', '· ') from a list line so only ONE marker shows. Safe on lines that
    have no marker (returns them unchanged)."""
    return _LEADING_LIST_MARKER_RE.sub("", text or "", count=1)


def _build_agenda_text(qa_pairs: list, report_structure: dict) -> str:
    """The agenda is the list of section titles the deck will cover: one
    bullet per finding (using `findings[i].page_title` if the brain
    provided one, otherwise a generic 'Finding N'), plus 'Key Takeaways'
    when we'll emit that slide too. Plain text, one entry per line — the
    renderer applies the agenda shape's font/size/color verbatim."""
    findings = report_structure.get("findings") or []
    lines: list[str] = []
    for i, _pair in enumerate(qa_pairs):
        raw = (findings[i].get("page_title") if i < len(findings) else "") or ""
        question = (_pair.get("human", {}).get("content")
                    if isinstance(_pair.get("human"), dict) else "") or ""
        page_title = _strip_emphasis(_finding_title(raw, question))
        lines.append(f"{i+1}. {page_title}")
    takeaways = report_structure.get("key_takeaways") or []
    if len(qa_pairs) >= 3 and takeaways:
        lines.append(f"{len(qa_pairs)+1}. Key Takeaways")
    return "\n".join(lines)




# ===========================================================================
# DESIGN-SPEC NATIVE RENDERER (version 3) — layout_plan.json driven
# ---------------------------------------------------------------------------
# Open the template as the base (theme/master/fonts/branding media inherited),
# PURGE every original slide, then ADD blank slides and place NATIVE, editable
# shapes (textboxes + pictures) at the coordinates the brain's layout_plan
# specifies. Re-place the template's branding image (by md5 match to the
# template media) but never clone the template's content-shaped boxes. Write an
# HTML QA preview alongside. Deliverable = native .pptx.
# ===========================================================================
EMU_PER_INCH = 914400


def _layout_plan_usable(plan) -> bool:
    """Mirror of brain.pptx_template_analyzer.layout_plan_usable."""
    if not isinstance(plan, dict):
        return False
    if int(plan.get("version") or 0) != 3:
        return False
    slides = plan.get("slides") or {}
    for k in ("cover", "agenda", "content", "takeaways"):
        if not isinstance(slides.get(k), dict):
            return False
    content = slides["content"]
    if not isinstance(content.get("title"), dict) or not isinstance(content.get("body"), dict):
        return False
    if not isinstance(content.get("chart"), dict):
        return False
    return True


def _template_media_index(template_path) -> dict:
    """Return {md5: (zip_member_name, bytes)} for every ppt/media/* image."""
    import hashlib
    import zipfile
    out: dict[str, tuple[str, bytes]] = {}
    try:
        with zipfile.ZipFile(str(template_path), "r") as z:
            for name in z.namelist():
                if name.startswith("ppt/media/"):
                    try:
                        b = z.read(name)
                    except Exception:
                        continue
                    out[hashlib.md5(b).hexdigest()] = (name, b)
    except Exception:
        return {}
    return out


def _pick_branding_bytes(media_index: dict, sid: str):
    """Pick the largest non-trivial template image as the branding logo.
    Returns (bytes, md5) or (None, None)."""
    if not media_index:
        return None, None
    candidates = [(md5, nm, b) for md5, (nm, b) in media_index.items() if len(b) > 512]
    if not candidates:
        candidates = [(md5, nm, b) for md5, (nm, b) in media_index.items()]
    if not candidates:
        return None, None
    candidates.sort(key=lambda t: len(t[2]), reverse=True)
    md5, nm, b = candidates[0]
    log_with_sid(sid, "info", "PPTX_BRANDING_PICK", member=nm, bytes=len(b), md5=md5[:10])
    return b, md5


def _blank_layout(prs):
    """Return the template layout with the fewest placeholders (ideally 0)."""
    best, best_n = None, None
    for layout in prs.slide_layouts:
        try:
            n = len(list(layout.placeholders))
        except Exception:
            n = 999
        if best is None or n < best_n:
            best, best_n = layout, n
    return best or prs.slide_layouts[0]


def _strip_placeholders(slide):
    """Remove placeholder shapes a blank layout still injected."""
    try:
        for shp in list(slide.shapes):
            if getattr(shp, "is_placeholder", False):
                el = shp._element
                parent = el.getparent()
                if parent is not None:
                    parent.remove(el)
    except Exception:
        pass


def _purge_all_slides(prs, sid: str) -> int:
    """Remove EVERY original slide, notesSlide-rel-safe. Returns count purged."""
    REL_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    try:
        sldIdLst = prs.slides._sldIdLst
    except Exception as e:
        log_with_sid(sid, "warning", f"PPTX_PURGE_NO_SLDIDLST: {e}")
        return 0
    originals = list(sldIdLst)
    purged = 0
    for sld_id in originals:
        rId = sld_id.attrib.get(REL_NS + "id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        try:
            sldIdLst.remove(sld_id)
            purged += 1
        except Exception:
            pass
    log_with_sid(sid, "info", "PPTX_PURGE_ALL_SLIDES", purged=purged,
                 remaining=len(prs.slides))
    return purged


def _align_enum(name):
    from pptx.enum.text import PP_ALIGN
    return {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT}.get((name or "left").lower(), PP_ALIGN.LEFT)


def _valign_enum(name):
    from pptx.enum.text import MSO_ANCHOR
    return {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM}.get((name or "top").lower(), MSO_ANCHOR.TOP)


def _chars_per_line(width_in: float, size_pt: float) -> int:
    if size_pt <= 0:
        size_pt = 12.0
    char_w_in = (0.50 * size_pt) / 72.0
    return max(1, int((width_in / char_w_in)))


def _fit_font_for_text(text: str, region: dict) -> float:
    """Shrink font until text plausibly fits the region height (>=9pt).
    Mirrors the checker's G4 capacity model so renderer + checker agree."""
    base = float(region.get("size_pt") or 14.0)
    w = float(region.get("w") or 6.0)
    h = float(region.get("h") or 3.0)
    n = len(text or "")
    size = base
    while size >= 9.0:
        cpl = _chars_per_line(w, size)
        line_h_in = (size * 1.2) / 72.0
        max_lines = max(1, int(h / line_h_in))
        if cpl * max_lines >= n:
            return size
        size -= 1.0
    return 9.0


def _adaptive_body_spacing(lines, size_pt, w, h):
    """Pick a line spacing (1.0-1.45) + per-bullet space-after so short body
    bullets distribute down the slot instead of bunching at the top, WITHOUT
    overflowing it (iteration 6). Returns (line_spacing, space_after_pt in pt).

    Estimates natural text height from the wrapped visual-line count at the
    already-fitted font size, then expands spacing only into genuine slack."""
    try:
        bullets = [ln for ln in lines if (ln or "").strip()]
        n = len(bullets)
        if n == 0 or size_pt <= 0 or h <= 0:
            return 1.0, 0.0
        one_line_in = size_pt * 1.2 / 72.0
        cpl = max(1, _chars_per_line(w, size_pt))
        vlines = 0
        for ln in bullets:
            chars = len(_strip_emphasis(ln))
            vlines += max(1, -(-chars // cpl))         # ceil division
        natural_in = vlines * one_line_in
        if natural_in <= 0:
            return 1.0, 0.0
        slack = h / natural_in
        if slack <= 1.15:
            return 1.0, 0.0                            # tight: leave it to shrink-fit
        spacing = min(1.45, max(1.0, slack * 0.85))
        # Never let expanded spacing overflow the slot.
        while spacing > 1.0 and vlines * one_line_in * spacing > h * 0.98:
            spacing -= 0.05
        spacing = round(max(1.0, spacing), 2)
        remaining_in = h - vlines * one_line_in * spacing
        space_after_pt = 0.0
        if remaining_in > 0 and n > 1:
            space_after_pt = min(10.0, (remaining_in / n) * 72.0 * 0.6)
        return spacing, round(max(0.0, space_after_pt), 1)
    except Exception:
        return 1.0, 0.0


def _add_text_region(slide, text: str, region: dict, *, sid: str,
                     shrink_to_fit: bool = False, line_spacing: float = None,
                     role: str = "", name: str = None, emphasis_color=None):
    """Add a native textbox at `region` (inches). word_wrap=True always.

    When `name` is given the shape is tagged with that name (the clone renderer
    passes a `PDC_INJ_*` sentinel so the checker can scope geometry checks to
    the content WE position). Native callers omit it → behaviour unchanged."""
    from pptx.util import Inches, Pt
    x, y = float(region["x"]), float(region["y"])
    w, h = float(region["w"]), float(region["h"])
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = _valign_enum(region.get("valign"))
    except Exception:
        pass
    base_pt = float(region.get("size_pt") or 14.0)
    # Titles: lower floor + NEVER ellipsis-truncated (iteration 6) — a <=4-word
    # title fits the header by shrinking the font alone. Body keeps its 9pt floor.
    floor_pt = 10.0 if role == "title" else 9.0
    size_pt = base_pt
    if shrink_to_fit:
        fitted = max(float(_fit_font_for_text(text or "", region)), floor_pt)
        if fitted < base_pt:
            log_with_sid(sid, "info", "PPTX_TEXT_OVERFLOW",
                         role=role, region_size_pt=base_pt, fitted_pt=fitted,
                         chars=len(text or ""))
        size_pt = fitted
        # Belt-and-suspenders so PowerPoint also shrinks on overflow.
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
    font_name = region.get("font") or "Calibri"
    color = _hex_to_rgb(region.get("color") or "374151")
    align = _align_enum(region.get("align"))
    lines = (text or "").split("\n")
    # Adaptive body/list spacing (iteration 6): when short bullets leave vertical
    # slack, expand line spacing + add space-after so they distribute down the
    # box instead of bunching at the top — without overflowing. Titles and
    # explicit line_spacing callers are untouched.
    eff_spacing = float(line_spacing) if line_spacing else None
    space_after_pt = 0.0
    if eff_spacing is None and shrink_to_fit and role in ("body", "list"):
        eff_spacing, space_after_pt = _adaptive_body_spacing(lines, size_pt, w, h)
    first = True
    for ln in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if eff_spacing:
            try:
                p.line_spacing = eff_spacing
            except Exception:
                pass
        if space_after_pt:
            try:
                p.space_after = Pt(space_after_pt)
            except Exception:
                pass
        # Titles are never emphasized; body/list lines may carry **bold** spans.
        _write_runs(p, ln, size_pt=size_pt, font_name=font_name,
                    color_rgb=color, bold=(role == "title"),
                    emphasis_color=(None if role == "title" else emphasis_color))
    return tb


def _add_branding(slide, brand_bytes, layout_plan, role: str, sid: str):
    """Place the branding image natively at the plan's branding region."""
    if not brand_bytes:
        return None
    from pptx.util import Inches
    branding = layout_plan.get("branding") or []
    placement = None
    for b in branding:
        roles = b.get("roles") or []
        if not roles or role in roles:
            placement = b
            break
    if placement is None and branding:
        placement = branding[0]
    if placement is None:
        return None
    x = float(placement.get("x_in") or 0.3)
    y = float(placement.get("y_in") or 0.3)
    w = float(placement.get("w_in") or 1.6)
    h = float(placement.get("h_in") or 0.6)
    try:
        slide.shapes.add_picture(io.BytesIO(brand_bytes), Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(h))
        return {"x": x, "y": y, "w": w, "h": h}
    except Exception as e:
        log_with_sid(sid, "warning", f"PPTX_BRANDING_PLACE_FAIL: {e}")
        return None


def _png_jpeg_size(b: bytes):
    """Return (w, h) for a PNG/JPEG, else None — no PIL dependency."""
    try:
        if b[:8] == b"\x89PNG\r\n\x1a\n":
            import struct
            w, h = struct.unpack(">II", b[16:24])
            return int(w), int(h)
        if b[:2] == b"\xff\xd8":
            import struct
            i = 2
            while i < len(b) - 9:
                if b[i] != 0xFF:
                    i += 1
                    continue
                marker = b[i + 1]
                if marker in (0xC0, 0xC1, 0xC2, 0xC3):
                    h, w = struct.unpack(">HH", b[i + 5:i + 9])
                    return int(w), int(h)
                seglen = struct.unpack(">H", b[i + 2:i + 4])[0]
                i += 2 + seglen
    except Exception:
        return None
    return None


def _place_chart_contain(slide, chart_b64, chart_region: dict, sid: str,
                         name: str | None = None):
    """Place the chart inside chart_region with 'contain' scaling — preserve
    aspect ratio, center, never exceed the region. Returns placed rect.

    `name` (optional) tags the placed picture (clone renderer passes
    `PDC_INJ_chart`); native callers omit it → behaviour unchanged."""
    if not chart_b64 or _is_plotly_html(chart_b64):
        return None
    from pptx.util import Inches
    try:
        img_bytes = base64.b64decode(chart_b64)
    except Exception as e:
        log_with_sid(sid, "warning", f"PPTX_CHART_DECODE_FAIL: {e}")
        return None
    wh = _png_jpeg_size(img_bytes)
    iw, ih = wh if wh else (1200, 800)
    rx, ry = float(chart_region["x"]), float(chart_region["y"])
    rw, rh = float(chart_region["w"]), float(chart_region["h"])
    if iw <= 0 or ih <= 0:
        iw, ih = 1200, 800
    img_ar = iw / ih
    region_ar = rw / rh
    if img_ar >= region_ar:
        draw_w = rw
        draw_h = rw / img_ar
    else:
        draw_h = rh
        draw_w = rh * img_ar
    off_x = rx + (rw - draw_w) / 2.0
    off_y = ry + (rh - draw_h) / 2.0
    try:
        pic = slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(off_x), Inches(off_y),
                                       width=Inches(draw_w), height=Inches(draw_h))
        if name:
            pic.name = name
        return {"x": off_x, "y": off_y, "w": draw_w, "h": draw_h, "img_ar": img_ar}
    except Exception as e:
        log_with_sid(sid, "warning", f"PPTX_CHART_PLACE_FAIL: {e}")
        return None


def _number_lines(text: str, bullet: str) -> str:
    lines = [ln for ln in (text or "").split("\n") if ln.strip()]
    if bullet == "dot":
        # Strip any enumerator/bullet the text already carries (e.g. the "N. "
        # _build_agenda_text prepends) before adding ours, so a fresh-region
        # agenda/takeaways list never shows TWO markers (PROBLEM 1).
        return "\n".join(f"• {_strip_leading_list_marker(ln)}" for ln in lines)
    return "\n".join(lines)



def _render_pptx_native(qa_pairs: list, report_structure: dict, sid: str,
                        template_path, layout_plan: dict) -> bytes:
    """Design-spec NATIVE renderer. Fresh editable deck on the template base
    using the layout_plan's per-slide-type regions. Order: cover -> agenda ->
    one content per finding -> takeaways (when >=3 findings + key_takeaways)."""
    from pptx import Presentation

    prs = Presentation(str(template_path))
    media_index = _template_media_index(template_path)
    brand_bytes, brand_md5 = _pick_branding_bytes(media_index, sid)

    _purge_all_slides(prs, sid)
    blank = _blank_layout(prs)

    slides_plan = layout_plan.get("slides") or {}
    cover_p = slides_plan.get("cover") or {}
    agenda_p = slides_plan.get("agenda") or {}
    content_p = slides_plan.get("content") or {}
    tk_p = slides_plan.get("takeaways") or {}

    preview_slides: list[dict] = []
    sw = float((layout_plan.get("slide_size_in") or {}).get("w") or 13.333)
    sh = float((layout_plan.get("slide_size_in") or {}).get("h") or 7.5)

    def _new_slide():
        s = prs.slides.add_slide(blank)
        _strip_placeholders(s)
        return s

    def _record(role, regions):
        preview_slides.append({"role": role, "regions": regions})

    # (a) Cover
    cover = _new_slide()
    regions = []
    br = _add_branding(cover, brand_bytes, layout_plan, "cover", sid)
    if br:
        regions.append({"kind": "branding", **br})
    title_txt = report_structure.get("report_title") or "Data Analysis Report"
    _add_text_region(cover, title_txt, cover_p.get("title") or {}, sid=sid, role="title")
    regions.append({"kind": "title", "text": title_txt, **(cover_p.get("title") or {})})
    sub = cover_p.get("subtitle")
    if sub:
        sub_txt = report_structure.get("executive_summary") or \
            f"Generated {datetime.utcnow().strftime('%Y-%m-%d')}"
        _add_text_region(cover, sub_txt, sub, sid=sid, shrink_to_fit=True, role="subtitle")
        regions.append({"kind": "subtitle", "text": sub_txt, **sub})
    _record("cover", regions)

    # (b) Agenda
    agenda = _new_slide()
    regions = []
    br = _add_branding(agenda, brand_bytes, layout_plan, "agenda", sid)
    if br:
        regions.append({"kind": "branding", **br})
    a_title = "Agenda"
    _add_text_region(agenda, a_title, agenda_p.get("title") or {}, sid=sid, role="title")
    regions.append({"kind": "title", "text": a_title, **(agenda_p.get("title") or {})})
    agenda_text = _build_agenda_text(qa_pairs, report_structure)
    a_list = agenda_p.get("list") or {}
    agenda_disp = _number_lines(agenda_text, a_list.get("bullet"))
    _add_text_region(agenda, agenda_disp, a_list, sid=sid, shrink_to_fit=True,
                     line_spacing=a_list.get("line_spacing"), role="list")
    regions.append({"kind": "list", "text": agenda_disp, **a_list})
    _record("agenda", regions)

    # (c) Content — one per finding
    findings = report_structure.get("findings") or []
    for i, pair in enumerate(qa_pairs):
        slide = _new_slide()
        regions = []
        br = _add_branding(slide, brand_bytes, layout_plan, "content", sid)
        if br:
            regions.append({"kind": "branding", **br})
        question = (pair["human"].get("content") or "") if isinstance(pair.get("human"), dict) else ""
        page_title = _finding_title(findings[i].get("page_title") if i < len(findings) else "", question)
        narrative = findings[i].get("narrative") if i < len(findings) else (pair["ai"].get("content") or "")
        _add_text_region(slide, page_title, content_p.get("title") or {}, sid=sid, role="title")
        regions.append({"kind": "title", "text": page_title, **(content_p.get("title") or {})})
        _add_text_region(slide, narrative, content_p.get("body") or {}, sid=sid,
                         shrink_to_fit=True, role="body")
        regions.append({"kind": "body", "text": narrative, **(content_p.get("body") or {})})
        chart_b64 = pair["ai"].get("_resolved_png") or pair["ai"].get("image_base64")
        placed = _place_chart_contain(slide, chart_b64, content_p.get("chart") or {}, sid)
        if placed:
            regions.append({"kind": "chart", **placed})
        _record("content", regions)

    # (d) Takeaways — when >=3 findings + key_takeaways present
    takeaways = report_structure.get("key_takeaways") or []
    if len(qa_pairs) >= 3 and takeaways:
        slide = _new_slide()
        regions = []
        br = _add_branding(slide, brand_bytes, layout_plan, "takeaways", sid)
        if br:
            regions.append({"kind": "branding", **br})
        t_title = "Key Takeaways"
        _add_text_region(slide, t_title, tk_p.get("title") or {}, sid=sid, role="title")
        regions.append({"kind": "title", "text": t_title, **(tk_p.get("title") or {})})
        t_list = tk_p.get("list") or {}
        tk_text = "\n".join(takeaways)
        tk_disp = _number_lines(tk_text, "dot")
        _add_text_region(slide, tk_disp, t_list, sid=sid, shrink_to_fit=True,
                         line_spacing=t_list.get("line_spacing"), role="list")
        regions.append({"kind": "list", "text": tk_disp, **t_list})
        _record("takeaways", regions)

    log_with_sid(sid, "info", "PPTX_RENDER_NATIVE_DONE",
                 slides=len(prs.slides), findings=len(qa_pairs),
                 takeaways=bool(len(qa_pairs) >= 3 and takeaways),
                 branding=bool(brand_bytes))

    try:
        _write_html_preview(preview_slides, sw, sh, sid)
    except Exception as e:
        log_with_sid(sid, "warning", f"PPTX_HTML_PREVIEW_FAIL: {e}")

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _html_preview_dir():
    from pathlib import Path as _P
    try:
        from settings import settings as _s
        base = _P(_s.DATA_ROOT) / "pptx_previews"
    except Exception:
        base = _P("/tmp/pptx_previews")
    base.mkdir(parents=True, exist_ok=True)
    return base


def _esc(s: str) -> str:
    return (str(s or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


def _write_html_preview(preview_slides: list, sw: float, sh: float, sid: str) -> str:
    """HTML QA preview — one <section> per slide at the same aspect ratio,
    absolute-positioned divs at the same inch->px scale. Plain HTML/CSS, no
    headless browser. QA ONLY — the deliverable is the native .pptx."""
    PX = 96
    out = [
        "<!doctype html><html><head><meta charset='utf-8'>",
        "<title>PPTX layout preview</title><style>",
        "body{background:#444;font-family:Arial,Helvetica,sans-serif;margin:0;padding:20px}",
        ".slide{position:relative;background:#fff;margin:0 auto 24px;"
        "box-shadow:0 2px 12px rgba(0,0,0,.4);overflow:hidden}",
        ".region{position:absolute;box-sizing:border-box;overflow:hidden}",
        ".title{font-weight:700}",
        ".branding{background:rgba(18,118,194,.15);border:1px dashed #1276c2}",
        ".chart{background:rgba(0,0,0,.05);border:1px solid #999;"
        "display:flex;align-items:center;justify-content:center;color:#666}",
        ".lbl{position:absolute;top:2px;left:4px;font-size:9px;color:#1276c2;opacity:.6}",
        "</style></head><body>",
        f"<p style='color:#fff'>QA preview — {len(preview_slides)} slides "
        f"({sw}in x {sh}in). Not the deliverable; the native .pptx is.</p>",
    ]
    for idx, sl in enumerate(preview_slides):
        out.append(f"<section class='slide' style='width:{sw*PX:.0f}px;height:{sh*PX:.0f}px'>")
        out.append(f"<div class='lbl'>#{idx+1} {_esc(sl.get('role'))}</div>")
        for r in sl.get("regions") or []:
            kind = r.get("kind")
            x = float(r.get("x") or 0); y = float(r.get("y") or 0)
            w = float(r.get("w") or 0); h = float(r.get("h") or 0)
            style = (f"left:{x*PX:.0f}px;top:{y*PX:.0f}px;"
                     f"width:{w*PX:.0f}px;height:{h*PX:.0f}px;")
            cls = "region"
            inner = ""
            if kind in ("title", "subtitle", "body", "list"):
                size = float(r.get("size_pt") or 14)
                color = "#" + (r.get("color") or "374151")
                align = r.get("align") or "left"
                style += (f"font-size:{size:.0f}px;color:{color};"
                          f"text-align:{align};white-space:pre-wrap;")
                if kind == "title":
                    cls += " title"
                inner = _esc((r.get("text") or ""))[:600]
            elif kind == "chart":
                cls += " chart"
                inner = "CHART"
            elif kind == "branding":
                cls += " branding"
            out.append(f"<div class='{cls}' style='{style}'>{inner}</div>")
        out.append("</section>")
    out.append("</body></html>")
    html = "".join(out)
    path = _html_preview_dir() / f"preview_{sid}.html"
    path.write_text(html, encoding="utf-8")
    log_with_sid(sid, "info", "PPTX_HTML_PREVIEW_WRITTEN", path=str(path),
                 slides=len(preview_slides))
    return str(path)


# ===========================================================================
# TEMPLATED CLONE RENDERER — clone the tenant's DESIGNED slides + inject
# ---------------------------------------------------------------------------
# Deep-clone each designed template slide (background, chrome, master/layout
# design and branding carry through unchanged), then apply the brain's v2
# build-plan labels: `drop` the template author's own sample shapes, inject our
# title / narrative / agenda / takeaways into the `replace:*` shapes (shrink to
# fit), and place the finding's chart in the analyzer-chosen region.
#
# Sentinel naming contract (so the checker can scope geometry checks to the
# content WE control — the template legitimately owns full-bleed backgrounds
# and edge-touching chrome that the native G1/G2/G3 checks were never meant to
# police):
#   PDC_INJ_<role>   fresh shape we positioned at OUR coordinates
#                    (chart, or a fallback textbox at a layout_plan region).
#                    Subject to in-bounds + capacity + overlap checks.
#   PDC_INJT_<role>  text injected into a template-owned shape — keeps the
#                    author's geometry; exempt from the in-bounds margin check
#                    (it inherits the template's placement) but still capacity-
#                    and overlap-checked against our other content.
#   (untagged)       `keep` chrome / background — original template names; the
#                    clone-aware checker ignores these for overlap/bounds.
# ===========================================================================
_REL_NS_BRACED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _remap_rels_in_el(el, rId_map: dict) -> None:
    """Rewrite r:embed / r:link / r:id references on a deep-copied element tree
    from the source slide's rIds to the destination slide's rIds."""
    if not rId_map:
        return
    for node in el.iter():
        for attr in list(node.attrib.keys()):
            if attr.startswith(_REL_NS_BRACED):
                new = rId_map.get(node.attrib[attr])
                if new:
                    node.set(attr, new)


def _clone_slide(prs, src_slide, sid: str):
    """Deep-clone `src_slide` into `prs` and return the new slide.

    python-pptx has no public slide-duplication API; we clone at the XML level:
      - create the new slide from the SOURCE slide's own layout (so the
        layout-level background / design is inherited, not a blank layout),
      - strip the placeholders add_slide injected,
      - re-create the source slide's relationships on the new slide part
        (image / chart / media), remapping rIds — SKIP notesSlide (it won't
        exist on the new slide and copying it corrupts the file) and slideLayout
        (add_slide already wired one),
      - deep-copy the slide-level background override (if any) and every shape,
        remapping their rId references to the new rIds.
    Never returns a partially-built slide: raises on failure so the caller can
    fall back to the native renderer (Article IV)."""
    import copy
    from pptx.oxml.ns import qn

    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    # Remove the placeholders the layout injected — we copy the source's own.
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    # Re-create relationships (remap rIds). add_slide already created the
    # slideLayout rel; notesSlide must be skipped.
    rId_map: dict[str, str] = {}
    for rId, rel in src_slide.part.rels.items():
        rt = rel.reltype
        if rt.endswith("slideLayout") or "notesSlide" in rt:
            continue
        try:
            if rel.is_external:
                new_rId = new_slide.part.relate_to(rel.target_ref, rt, is_external=True)
            else:
                new_rId = new_slide.part.relate_to(rel.target_part, rt)
            rId_map[rId] = new_rId
        except Exception as e:
            log_with_sid(sid, "warning", f"PPTX_CLONE_REL_FAIL rId={rId}: {e}")

    # Copy the slide-level background override, if the source overrides it.
    src_cSld = src_slide._element.find(qn("p:cSld"))
    new_cSld = new_slide._element.find(qn("p:cSld"))
    if src_cSld is not None and new_cSld is not None:
        src_bg = src_cSld.find(qn("p:bg"))
        if src_bg is not None:
            bg_copy = copy.deepcopy(src_bg)
            _remap_rels_in_el(bg_copy, rId_map)
            new_cSld.insert(0, bg_copy)  # bg must be the first child of cSld

    # Deep-copy every shape element from the source slide.
    spTree = new_slide.shapes._spTree
    for shp in src_slide.shapes:
        el = copy.deepcopy(shp._element)
        _remap_rels_in_el(el, rId_map)
        spTree.append(el)
    return new_slide


def _set_shape_rect(shape, region: dict) -> None:
    """Resize/move a cloned shape to `region` (inches). Used for the content
    body shape when a chart shares the slide, so narrative + chart never
    overlap (the template's single content box can't hold both)."""
    from pptx.util import Inches
    try:
        shape.left = Inches(float(region["x"]))
        shape.top = Inches(float(region["y"]))
        shape.width = Inches(float(region["w"]))
        shape.height = Inches(float(region["h"]))
    except Exception:
        pass


def _truncate_to_fit(text: str, region: dict, pt: float) -> str:
    """Truncate `text` at a word boundary with an ellipsis so it fits `region`
    (w×h inches) at `pt`. Same char-box model as `_fit_font_for_text`. Used as a
    last resort for titles that still overflow at the minimum font size."""
    txt = (text or "").strip()
    if not txt:
        return txt
    w = float(region.get("w") or 4.0)
    h = float(region.get("h") or 1.0)
    size_in = max(float(pt), 1.0) / 72.0
    cpl = max(1, int(w / (0.6 * size_in)))
    lines_avail = max(1, int(h / (1.2 * size_in)))
    budget = cpl * lines_avail
    if len(txt) <= budget:
        return txt
    cut = txt[:max(1, budget - 1)].rstrip()
    if " " in cut:
        cut = cut.rsplit(" ", 1)[0].rstrip()
    return (cut + "…") if cut else (txt[:1] + "…")


def _fill_plain(tf, lines, *, size_pt, font_name, color_rgb, bold, emphasis_color=None) -> bool:
    """Clear the frame and write `lines` as plain runs (no bullet structure).

    When `emphasis_color` is given, **...** spans within a line render bold in
    that color (markers stripped); otherwise the markers are stripped to plain."""
    tf.clear()
    for idx, ln in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        _write_runs(p, ln, size_pt=size_pt, font_name=font_name,
                    color_rgb=color_rgb, bold=bold, emphasis_color=emphasis_color)
    return True


def _list_native_bullet(tf) -> bool:
    """True when the text frame's list resolves to a NATIVE bullet / auto-number
    that PowerPoint will actually render — checked across ALL paragraphs' own
    `<a:pPr>` AND the shape `txBody`'s `<a:lstStyle>` level properties (not just
    the first paragraph). A paragraph/level that carries an explicit `<a:buNone>`
    is treated as having no marker. Inherited placeholder/master list styles are
    NOT resolved here — when nothing explicit is found we return False so the
    caller forces a single deterministic marker (PROBLEM 1: inherited markers
    were invisible to the old `pPr0`-only check and stacked under our "N. ")."""
    try:
        for p in tf.paragraphs:
            pPr = p._p.find(qn("a:pPr"))
            if pPr is None:
                continue
            if pPr.find(qn("a:buNone")) is not None:
                continue
            if (pPr.find(qn("a:buChar")) is not None
                    or pPr.find(qn("a:buAutoNum")) is not None):
                return True
        lstStyle = tf._txBody.find(qn("a:lstStyle"))
        if lstStyle is not None:
            for lvl in lstStyle:
                if lvl.find(qn("a:buNone")) is not None:
                    continue
                if (lvl.find(qn("a:buChar")) is not None
                        or lvl.find(qn("a:buAutoNum")) is not None):
                    return True
    except Exception:
        return False
    return False


def _suppress_paragraph_bullet(p) -> None:
    """Force `<a:buNone>` on a paragraph's `<a:pPr>` (removing any buChar /
    buAutoNum first) so NO native/inherited marker renders — used when WE supply
    the explicit enumerator, guaranteeing exactly one marker. buNone is inserted
    in its schema-correct slot (before tabLst/defRPr/extLst). Never raises."""
    try:
        pPr = p._p.get_or_add_pPr()
        for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
            for el in pPr.findall(qn(tag)):
                pPr.remove(el)
        buNone = pPr.makeelement(qn("a:buNone"), {})
        ref = None
        for tag in ("a:tabLst", "a:defRPr", "a:extLst"):
            ref = pPr.find(qn(tag))
            if ref is not None:
                break
        if ref is not None:
            ref.addprevious(buNone)
        else:
            pPr.append(buNone)
    except Exception:
        pass


def _fill_keep_bullets(tf, lines, *, size_pt, font_name, color_rgb, bold, emphasis_color=None,
                       own_marker: bool = False) -> bool:
    """Fill `lines` into `tf` while REUSING the template paragraphs' bullet /
    number / indent formatting (the paragraph-level <a:pPr>), so a list keeps the
    template's exact list structure instead of being flattened to plain text.

    `_Paragraph.clear()` removes runs but preserves <a:pPr>, so we clear each
    paragraph and write a fresh run; overflow lines clone the last (formatted)
    paragraph element. Returns False when there is no paragraph to use as a style
    template (the caller then falls back to plain runs).

    `own_marker=True` (agenda / takeaways — WE prepend "N." / "•"): guarantee
    EXACTLY ONE marker. If a native bullet/number is confirmed, strip our
    explicit one (template's shows); otherwise force `<a:buNone>` on every
    injected paragraph and keep ours — so an inherited/master marker can never
    stack under it. `own_marker=False` (content body prose — no explicit marker):
    legacy behavior, reuse the template's list formatting untouched."""
    from pptx.util import Pt
    import copy
    try:
        if not tf.paragraphs:
            return False
        if own_marker:
            # Agenda / takeaways: we own the enumerator → enforce one marker.
            native_marker = _list_native_bullet(tf)
            if native_marker:
                lines = [_strip_leading_list_marker(ln) for ln in lines]
        else:
            # Content body (prose): preserve old pPr0-only behavior — strip an
            # explicit leading marker only when the first paragraph carries an
            # explicit native one (a no-op for ordinary prose). Never force buNone
            # so legitimate template body bullets are untouched.
            native_marker = True
            pPr0 = tf.paragraphs[0]._p.find(qn("a:pPr"))
            if pPr0 is not None and (
                pPr0.find(qn("a:buChar")) is not None
                or pPr0.find(qn("a:buAutoNum")) is not None
            ):
                lines = [_strip_leading_list_marker(ln) for ln in lines]
        txBody = tf._txBody
        last_el = tf.paragraphs[-1]._p
        while len(tf.paragraphs) < len(lines):
            txBody.append(copy.deepcopy(last_el))
        for i, p in enumerate(list(tf.paragraphs)):
            if i < len(lines):
                p.clear()                       # drops runs, KEEPS pPr (bullet)
                _write_runs(p, lines[i], size_pt=size_pt, font_name=font_name,
                            color_rgb=color_rgb, bold=bold,
                            emphasis_color=emphasis_color)
                # own_marker + no confirmed native marker → suppress any
                # inherited/master bullet so only our "N." / "•" shows.
                if own_marker and not native_marker:
                    _suppress_paragraph_bullet(p)
            else:
                p._p.getparent().remove(p._p)   # drop surplus template paragraphs
        return True
    except Exception:
        return False


def _inject_text(shape, text: str, *, role: str, sid: str,
                 shrink_region: dict | None = None, fallback_style: dict | None = None,
                 emphasis_color=None) -> bool:
    """Replace a cloned template shape's text with `text`, preserving the
    shape's own first-run formatting where possible (Article XI — reuse the
    template's font/size/color); only override the size when shrinking to fit.
    Titles collapse to a single run so the checker's "appears exactly once"
    assertion holds.

    Returns True when the text was injected. Returns False when the target
    shape cannot hold text (e.g. a table `graphicFrame` or a picture — which is
    exactly what some build plans label `replace:body`); the caller then drops
    that sample shape and renders the text in a fresh region instead."""
    from pptx.util import Pt
    try:
        tf = shape.text_frame
    except Exception:
        return False
    tf.word_wrap = True

    # Capture the template's first-run formatting before we clear it.
    tmpl_name = tmpl_bold = tmpl_color = None
    tmpl_size_pt = None
    for para in tf.paragraphs:
        if para.runs:
            f = para.runs[0].font
            tmpl_name = f.name
            tmpl_bold = f.bold
            try:
                if f.size is not None:
                    tmpl_size_pt = f.size.pt
            except Exception:
                tmpl_size_pt = None
            try:
                if f.color is not None and f.color.type is not None:
                    tmpl_color = f.color.rgb
            except Exception:
                tmpl_color = None
            break

    fs = fallback_style or {}
    font_name = tmpl_name or fs.get("font") or "Calibri"
    if tmpl_color is not None:
        color_rgb = tmpl_color
    else:
        color_rgb = _hex_to_rgb(fs.get("color_hex") or "374151")
    base_pt = tmpl_size_pt or float(fs.get("size_pt") or (26.0 if role == "title" else 14.0))

    # Fit the text to the slot. Titles in particular MUST never overflow their
    # header box (iteration 4): shrink toward a floor, then truncate with an
    # ellipsis if even the floor will not fit.
    size_pt = base_pt
    # Titles: lower floor + NEVER truncated with an ellipsis (iteration 6).
    floor_pt = 10.0 if role == "title" else 9.0
    region = None
    if shrink_region and shrink_region.get("w") and shrink_region.get("h"):
        region = {"size_pt": base_pt, "w": shrink_region["w"], "h": shrink_region["h"]}
        fitted = max(float(_fit_font_for_text(text or "", region)), floor_pt)
        if fitted < base_pt:
            log_with_sid(sid, "info", "PPTX_TEXT_OVERFLOW", role=role,
                         region_size_pt=base_pt, fitted_pt=fitted, chars=len(text or ""))
            size_pt = fitted

    if role == "title":
        bold = True
    elif tmpl_bold is not None:
        bold = bool(tmpl_bold)
    else:
        bold = False

    # Titles are NEVER ellipsis-truncated (iteration 6) — shrink-to-fit + the
    # PowerPoint autofit below keep a <=4-word title fully visible on the header.
    text_out = text or ""

    # Belt-and-suspenders: ask PowerPoint to shrink text on overflow so the slot
    # is never exceeded even if the heuristic under-estimates.
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    # Titles are never emphasized; body/list inject may carry the theme accent.
    emph = None if role == "title" else emphasis_color
    lines = text_out.split("\n")
    if role in ("body", "list") and len(lines) > 1:
        # Preserve the template list's bullet / numbering / indent structure.
        # role=="list" (agenda / takeaways) → WE own the enumerator, so enforce
        # exactly one marker; role=="body" (prose) keeps the template's bullets.
        if not _fill_keep_bullets(tf, lines, size_pt=size_pt, font_name=font_name,
                                  color_rgb=color_rgb, bold=bold, emphasis_color=emph,
                                  own_marker=(role == "list")):
            _fill_plain(tf, lines, size_pt=size_pt, font_name=font_name,
                        color_rgb=color_rgb, bold=bold, emphasis_color=emph)
    else:
        _fill_plain(tf, lines, size_pt=size_pt, font_name=font_name,
                    color_rgb=color_rgb, bold=bold, emphasis_color=emph)
    return True


# ---------------------------------------------------------------------------
# Placement validation (iteration 3) — our content must NEVER land on a logo
# or off-canvas. A spec-labeled replace:* shape is used in place only when it
# is a genuine, on-canvas content slot clear of the template's logos; otherwise
# the template shape is left untouched (it is part of the design) and our text
# is rendered in a computed SAFE region instead — or, on a logo-only cover,
# nothing is added at all.
# ---------------------------------------------------------------------------
def _emu_to_in(v):
    from pptx.util import Emu
    try:
        return round(Emu(v).inches, 3) if v is not None else None
    except Exception:
        return None


def _shape_rect_in(shape):
    """(x, y, w, h) in inches, or None when the shape has no explicit geometry
    (a placeholder inheriting its layout position — treated as on-canvas)."""
    x, y = _emu_to_in(shape.left), _emu_to_in(shape.top)
    w, h = _emu_to_in(shape.width), _emu_to_in(shape.height)
    if None in (x, y, w, h):
        return None
    return (x, y, w, h)


def _can_hold_text(shape) -> bool:
    try:
        _ = shape.text_frame
        return True
    except Exception:
        return False


def _norm_txt(s: str) -> str:
    return " ".join((s or "").split()).strip().lower()


def _kept_text_matches(slide, label_for: dict, target: str) -> bool:
    """True when a KEEP shape on the slide already displays `target` — so we must
    not inject a duplicate title (e.g. an 'Agenda' wordmark that is kept chrome).
    Conservative to avoid false positives on long content titles: exact
    normalized match, or — only for short labels (<=24 chars) — a containment
    match either way."""
    tnorm = _norm_txt(target)
    if not tnorm:
        return False
    for shp in slide.shapes:
        try:
            sid_int = int(shp.shape_id)
        except Exception:
            sid_int = None
        if label_for.get(sid_int, "keep") != "keep":
            continue
        if not _can_hold_text(shp):
            continue
        try:
            txt = _norm_txt(shp.text_frame.text)
        except Exception:
            continue
        if not txt:
            continue
        if txt == tnorm:
            return True
        if len(tnorm) <= 24 and (tnorm in txt or txt in tnorm):
            return True
    return False


def _is_picture(shape) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except Exception:
        pass
    try:
        return shape._element.tag.endswith("}pic")
    except Exception:
        return False


def _rect_overlap_frac(a, b) -> float:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    ov = ix * iy
    amin = min(aw * ah, bw * bh)
    return (ov / amin) if amin > 0 else 0.0


def _protected_logos(slide, label_for: dict, sw: float, sh: float) -> list:
    """Rects of KEEP pictures small enough to be logos — our text must never
    land on them. Large kept banners / full-bleed backgrounds are NOT logos
    (a title may legitimately sit on a banner)."""
    out = []
    area = (sw or 13.333) * (sh or 7.5)
    for shp in slide.shapes:
        try:
            sid_int = int(shp.shape_id)
        except Exception:
            sid_int = None
        if label_for.get(sid_int, "keep") != "keep":
            continue
        if not _is_picture(shp):
            continue
        r = _shape_rect_in(shp)
        if not r or r[2] <= 0 or r[3] <= 0:
            continue
        _, _, w, h = r
        if (w < 0.40 * sw and h < 0.25 * sh) or (w * h < 0.12 * area):
            out.append(r)
    return out


def _on_canvas(rect, sw: float, sh: float, tol: float = 0.05) -> bool:
    if rect is None:
        return True
    x, y, w, h = rect
    return (x >= -tol and y >= -tol
            and x + w <= sw + tol and y + h <= sh + tol)


def _position_ok(shape, sw: float, sh: float, logos: list) -> bool:
    """True when `shape` is on-canvas AND clear of every protected logo
    (<= 10% of the smaller area). A shape failing this must not receive our
    text at its template geometry."""
    rect = _shape_rect_in(shape)
    if not _on_canvas(rect, sw, sh):
        return False
    if rect is not None:
        for lg in logos:
            if _rect_overlap_frac(rect, lg) > 0.10:
                return False
    return True


def _region_safe(region, sw: float, sh: float, logos: list) -> bool:
    if not (isinstance(region, dict) and all(k in region for k in ("x", "y", "w", "h"))):
        return False
    rect = (float(region["x"]), float(region["y"]),
            float(region["w"]), float(region["h"]))
    if not _on_canvas(rect, sw, sh):
        return False
    return all(_rect_overlap_frac(rect, lg) <= 0.10 for lg in logos)


def _computed_region(role: str, sw: float, sh: float, logos: list) -> dict:
    """A generic safe region for a role, clear of protected logos. Title sits
    in the top header band (trimmed away from any corner logo); body/list fill
    the main area below the header."""
    M, HEADER_TOP, HEADER_H, HEADER_BOTTOM = 0.5, 0.3, 0.8, 1.4
    if role == "title":
        x, w = M, sw - 2 * M
        for (lx, ly, lw, lh) in logos:
            if ly < HEADER_TOP + HEADER_H and lx > sw / 2:   # logo in right header
                w = min(w, lx - x - 0.15)
        return {"x": round(x, 2), "y": HEADER_TOP, "w": round(max(2.0, w), 2), "h": HEADER_H}
    return {"x": M, "y": HEADER_BOTTOM,
            "w": round(sw - 2 * M, 2), "h": round(sh - HEADER_BOTTOM - M, 2)}


def _safe_region(role: str, sw: float, sh: float, logos: list, plan_region=None) -> dict:
    """Prefer the brain's layout_plan region when it is on-canvas and clear of
    logos; otherwise compute a safe zone."""
    if _region_safe(plan_region, sw, sh, logos):
        return plan_region
    return _computed_region(role, sw, sh, logos)


def _render_pptx_templated_clone(qa_pairs: list, report_structure: dict, sid: str,
                                 template_path, spec: dict, layout_plan: dict | None = None) -> bytes:
    """Clone the tenant's DESIGNED slides and overlay our content. Deck order:
    cover -> agenda (when the v2 deck names one) -> one content per finding ->
    takeaways (when >=3 findings + key_takeaways). The content body shape is
    resized to the validated layout_plan region when a chart shares the slide
    so narrative + chart never overlap; otherwise it keeps the template's
    geometry."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(template_path))
    deck = spec.get("deck") or {}
    cover_i = deck.get("cover_slide_index")
    agenda_i = deck.get("agenda_slide_index")
    content_i = deck.get("content_slide_index")
    slides_spec = spec.get("slides") or {}

    sw_in = round(Emu(prs.slide_width).inches, 2)
    sh_in = round(Emu(prs.slide_height).inches, 2)

    src = list(prs.slides)

    def _src(i):
        return src[i] if isinstance(i, int) and 0 <= i < len(src) else None

    cover_src = _src(cover_i)
    content_src = _src(content_i)
    agenda_src = _src(agenda_i)
    if cover_src is None or content_src is None:
        raise ValueError("clone: cover/content source slide missing from template")

    # Capture the original sldId elements BEFORE cloning so we can purge exactly
    # the originals afterwards (clones are appended; indexes stay valid while we
    # clone).
    original_sldIds = list(prs.slides._sldIdLst)

    lp_content = ((layout_plan or {}).get("slides") or {}).get("content") or {}

    def _content_regions():
        """Return (body_region, chart_region) for a content slide carrying a
        chart. Prefer the validated layout_plan; else split the slide."""
        body = lp_content.get("body")
        chart = lp_content.get("chart")
        if isinstance(body, dict) and isinstance(chart, dict):
            return body, chart
        sw = round(Emu(prs.slide_width).inches, 2)
        sh = round(Emu(prs.slide_height).inches, 2)
        m, top = 0.3, 1.6
        body = {"x": m, "y": top, "w": round((sw - 2 * m) * 0.44, 2), "h": round(sh - top - m, 2)}
        chart = {"x": round(m + (sw - 2 * m) * 0.48, 2), "y": top,
                 "w": round((sw - 2 * m) * 0.52, 2), "h": round(sh - top - m, 2), "fit": "contain"}
        return body, chart

    def _apply(slide, spec_slide, logos, *, page_title=None, body_text=None,
               agenda_text=None, body_region=None, accent_rgb=None):
        """Apply the v2 shape labels to a cloned slide + inject our content.

        A `replace:*` shape is only written in place when it is a VALID slot:
        it can hold text AND sits on-canvas clear of every protected logo. If
        it cannot hold text (a table / picture sample) it is dropped; if it is
        a text shape in a bad position (off-canvas, or over a logo — e.g. the
        analyzer's backstop mislabeling a full-bleed decorative rectangle on a
        logo-only cover) it is LEFT UNTOUCHED as part of the design, and the
        caller renders our text in a computed safe region instead.

        Returns a dict of which roles were placed in place (so the caller adds
        a fresh fallback shape for any it did not)."""
        shapes_spec = (spec_slide or {}).get("shapes") or []
        label_for: dict[int, str] = {}
        style_for: dict[int, dict] = {}
        for entry in shapes_spec:
            sidv = entry.get("shape_id")
            if isinstance(sidv, int):
                label_for[sidv] = entry.get("label") or "keep"
                style_for[sidv] = entry.get("text_style") or {}
        placed = {"title": False, "body": False, "agenda": False}
        # No duplicate title: if a KEPT chrome shape already displays this title
        # (e.g. an "Agenda" wordmark), don't inject another one and don't add a
        # fresh fallback either.
        if page_title is not None and _kept_text_matches(slide, label_for, page_title):
            placed["title"] = True
            log_with_sid(sid, "info", "PPTX_TITLE_DEDUP", role="title")
        for shp in list(slide.shapes):
            try:
                sid_int = int(shp.shape_id)
            except Exception:
                continue
            label = label_for.get(sid_int, "keep")
            if label == "drop":
                shp._element.getparent().remove(shp._element)
            elif label == "replace:title" and page_title is not None and not placed["title"]:
                if not _can_hold_text(shp):
                    shp._element.getparent().remove(shp._element)   # table/pic sample
                elif _position_ok(shp, sw_in, sh_in, logos):
                    r = _shape_rect_in(shp)
                    sr = {"w": r[2], "h": r[3]} if r else None
                    if _inject_text(shp, page_title, role="title", sid=sid,
                                    shrink_region=sr, fallback_style=style_for.get(sid_int)):
                        shp.name = "PDC_INJT_title"
                        placed["title"] = True
                # else: bad position -> leave the template shape untouched; the
                # caller renders the title in a safe region (or nothing on a
                # logo-only cover).
            elif label == "replace:body" and body_text is not None:
                if body_region is not None:
                    # A chart shares this slide: never reuse the template body
                    # shape (it would collide with the chart, and is often a
                    # decorative header band). Leave it as design; the caller
                    # adds a fresh PDC_INJ_body in the validated body region.
                    pass
                elif not _can_hold_text(shp):
                    shp._element.getparent().remove(shp._element)   # table/pic sample
                elif _position_ok(shp, sw_in, sh_in, logos):
                    sr = {"w": round(Emu(shp.width).inches, 2) if shp.width else None,
                          "h": round(Emu(shp.height).inches, 2) if shp.height else None}
                    if _inject_text(shp, body_text, role="body", sid=sid,
                                    shrink_region=sr, fallback_style=style_for.get(sid_int),
                                    emphasis_color=accent_rgb):
                        shp.name = "PDC_INJT_body"             # keep template geometry
                        placed["body"] = True
                # else: bad position -> leave untouched; caller adds fresh region.
            elif label == "replace:agenda" and agenda_text is not None:
                if not _can_hold_text(shp):
                    shp._element.getparent().remove(shp._element)
                elif _position_ok(shp, sw_in, sh_in, logos):
                    sr = {"w": round(Emu(shp.width).inches, 2) if shp.width else None,
                          "h": round(Emu(shp.height).inches, 2) if shp.height else None}
                    if _inject_text(shp, agenda_text, role="list", sid=sid,
                                    shrink_region=sr, fallback_style=style_for.get(sid_int),
                                    emphasis_color=accent_rgb):
                        shp.name = "PDC_INJT_list"
                        placed["agenda"] = True
                # else: bad position -> leave untouched; caller adds fresh region.
            # keep -> leave untouched
        return placed

    def _logos_for(slide, spec_slide):
        lf = {}
        for e in (spec_slide or {}).get("shapes") or []:
            if isinstance(e.get("shape_id"), int):
                lf[e["shape_id"]] = e.get("label") or "keep"
        return _protected_logos(slide, lf, sw_in, sh_in)

    def _fallback_region(slide_type, role):
        sl = ((layout_plan or {}).get("slides") or {}).get(slide_type) or {}
        r = sl.get(role)
        return r if isinstance(r, dict) else None

    findings = report_structure.get("findings") or []
    takeaways = report_structure.get("key_takeaways") or []
    want_takeaways = len(qa_pairs) >= 3 and bool(takeaways)

    # Theme accent for **bold** emphasis spans (iteration 5). Computed once;
    # threaded into body/agenda/takeaways injects only (never titles, never cover).
    accent_rgb = _hex_to_rgb(((layout_plan or {}).get("palette") or {}).get("accent")
                             or ((spec or {}).get("theme_colors") or {}).get("accent")
                             or "1276C2")

    # (a) Cover — inject a title ONLY into a genuine, on-canvas title slot.
    #     A logo-only cover (no valid title slot — the analyzer's backstop may
    #     have mislabeled a full-bleed rectangle over the logo) is cloned
    #     untouched so it matches the template exactly. We never write over the
    #     logo, and we add NO fresh fallback on the cover.
    cover = _clone_slide(prs, cover_src, sid)
    cover_spec = slides_spec.get(str(cover_i))
    title_txt = report_structure.get("report_title") or "Data Analysis Report"
    _apply(cover, cover_spec, _logos_for(cover, cover_spec), page_title=title_txt)

    # (b) Agenda — only when the template named an agenda slide.
    if agenda_src is not None and str(agenda_i) in slides_spec:
        agenda = _clone_slide(prs, agenda_src, sid)
        agenda_spec = slides_spec.get(str(agenda_i))
        a_logos = _logos_for(agenda, agenda_spec)
        agenda_text = _build_agenda_text(qa_pairs, report_structure)
        placed = _apply(agenda, agenda_spec, a_logos,
                        page_title="Agenda", agenda_text=agenda_text, accent_rgb=accent_rgb)
        if not placed["title"]:
            reg = _safe_region("title", sw_in, sh_in, a_logos,
                               _fallback_region("agenda", "title"))
            _add_text_region(agenda, "Agenda", reg, sid=sid, shrink_to_fit=True,
                             role="title", name="PDC_INJ_title")
        if not placed["agenda"]:
            reg = _safe_region("list", sw_in, sh_in, a_logos, _fallback_region("agenda", "list"))
            _add_text_region(agenda, _number_lines(agenda_text, reg.get("bullet")), reg,
                             sid=sid, shrink_to_fit=True, role="list", name="PDC_INJ_list",
                             emphasis_color=accent_rgb)

    # (c) Content — one clone per finding.
    content_spec = slides_spec.get(str(content_i))
    body_region, chart_region = _content_regions()
    for i, pair in enumerate(qa_pairs):
        slide = _clone_slide(prs, content_src, sid)
        c_logos = _logos_for(slide, content_spec)
        question = (pair["human"].get("content") or "") if isinstance(pair.get("human"), dict) else ""
        page_title = _finding_title(findings[i].get("page_title") if i < len(findings) else "", question)
        narrative = findings[i].get("narrative") if i < len(findings) else (pair["ai"].get("content") or "")
        chart_b64 = pair["ai"].get("_resolved_png") or pair["ai"].get("image_base64")
        has_chart = bool(chart_b64) and not _is_plotly_html(chart_b64)
        placed = _apply(slide, content_spec, c_logos, page_title=page_title, body_text=narrative,
                        body_region=(body_region if has_chart else None), accent_rgb=accent_rgb)
        if not placed["title"]:
            reg = _safe_region("title", sw_in, sh_in, c_logos, _fallback_region("content", "title"))
            _add_text_region(slide, page_title, reg, sid=sid, shrink_to_fit=True,
                             role="title", name="PDC_INJ_title")
        if not placed["body"]:
            plan_body = body_region if has_chart else _fallback_region("content", "body")
            reg = _safe_region("body", sw_in, sh_in, c_logos, plan_body)
            # Fresh content body: render each narrative line as a dot bullet.
            body_disp = _number_lines(narrative, "dot")
            _add_text_region(slide, body_disp, reg, sid=sid, shrink_to_fit=True,
                             role="body", name="PDC_INJ_body", emphasis_color=accent_rgb)
        if has_chart:
            _place_chart_contain(slide, chart_b64, chart_region, sid, name="PDC_INJ_chart")

    # (d) Takeaways — reuse the content clone's chrome; no chart.
    if want_takeaways:
        slide = _clone_slide(prs, content_src, sid)
        t_logos = _logos_for(slide, content_spec)
        tk_text = _number_lines("\n".join(takeaways), "dot")
        placed = _apply(slide, content_spec, t_logos, page_title="Key Takeaways",
                        body_text=tk_text, accent_rgb=accent_rgb)
        if not placed["title"]:
            reg = _safe_region("title", sw_in, sh_in, t_logos, _fallback_region("content", "title"))
            _add_text_region(slide, "Key Takeaways", reg, sid=sid, shrink_to_fit=True,
                             role="title", name="PDC_INJ_title")
        if not placed["body"]:
            reg = _safe_region("body", sw_in, sh_in, t_logos,
                               _fallback_region("takeaways", "list") or _fallback_region("content", "body"))
            _add_text_region(slide, tk_text, reg, sid=sid, shrink_to_fit=True,
                             role="list", name="PDC_INJ_list", emphasis_color=accent_rgb)

    # Purge the original (now-cloned) template slides.
    sldIdLst = prs.slides._sldIdLst
    purged = 0
    for sld_id in original_sldIds:
        rId = sld_id.attrib.get(_REL_NS_BRACED + "id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        try:
            sldIdLst.remove(sld_id)
            purged += 1
        except Exception:
            pass

    log_with_sid(sid, "info", "PPTX_RENDER_CLONE_DONE",
                 slides=len(prs.slides), purged=purged, findings=len(qa_pairs),
                 agenda=bool(agenda_src is not None and str(agenda_i) in slides_spec),
                 takeaways=want_takeaways)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _render_pdf(qa_pairs: list, report_structure: dict, sid: str) -> bytes:
    """Minimal PDF rendering — same brand palette + DejaVu fonts as B2C.

    Charts are embedded as PNGs (already resolved by `_build_qa_pairs` via
    kaleido for Plotly). Tables are rendered as ReportLab tables.
    """
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image as RLImage,
        Table as RLTable, TableStyle,
    )
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.colors import HexColor, white
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from pathlib import Path as _P

    # Register DejaVu Sans (Georgian support) if present
    _here = _P(__file__).resolve().parent.parent
    fonts_dir = _here / "static" / "fonts"
    fn = "Helvetica"
    fn_bold = "Helvetica-Bold"
    try:
        reg = fonts_dir / "DejaVuSans.ttf"
        bold = fonts_dir / "DejaVuSans-Bold.ttf"
        if reg.exists() and bold.exists():
            pdfmetrics.registerFont(TTFont("DejaVu", str(reg)))
            pdfmetrics.registerFont(TTFont("DejaVu-Bold", str(bold)))
            fn, fn_bold = "DejaVu", "DejaVu-Bold"
    except Exception:
        pass

    title_style = ParagraphStyle("T", fontName=fn_bold, fontSize=24, leading=30, textColor=HexColor("#001E44"))
    h2 = ParagraphStyle("H2", fontName=fn_bold, fontSize=16, leading=22, spaceBefore=12, textColor=HexColor("#001E44"))
    body = ParagraphStyle("B", fontName=fn, fontSize=11, leading=15, textColor=HexColor("#374151"))
    meta = ParagraphStyle("M", fontName=fn, fontSize=9, leading=12, textColor=HexColor("#9ca3af"))

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, leftMargin=40, rightMargin=40, topMargin=40, bottomMargin=40)
    flow = []

    flow.append(Paragraph(report_structure.get("report_title") or "Data Analysis Report", title_style))
    flow.append(Spacer(1, 12))
    es = report_structure.get("executive_summary") or ""
    if es:
        flow.append(Paragraph(es, body))
        flow.append(Spacer(1, 12))

    findings = report_structure.get("findings") or []
    for i, pair in enumerate(qa_pairs):
        flow.append(PageBreak())
        question = (pair["human"].get("content") or "") if isinstance(pair.get("human"), dict) else ""
        page_title = _finding_title(findings[i].get("page_title") if i < len(findings) else "", question)
        narrative = findings[i].get("narrative") if i < len(findings) else (pair["ai"].get("content") or "")
        flow.append(Paragraph(_strip_emphasis(page_title), h2))
        flow.append(Paragraph(_strip_emphasis(narrative), body))
        flow.append(Spacer(1, 12))

        # Chart
        chart_b64 = pair["ai"].get("_resolved_png") or pair["ai"].get("image_base64")
        if chart_b64 and not _is_plotly_html(chart_b64):
            try:
                img_bytes = base64.b64decode(chart_b64)
                img = RLImage(io.BytesIO(img_bytes), width=480, height=300)
                flow.append(img)
            except Exception as e:
                log_with_sid(sid, "warning", f"PDF_IMG_FAIL: {e}")

        # Table
        table = pair["ai"].get("table")
        if table and not chart_b64:
            cols = table.get("columns") or []
            rows = (table.get("rows") or [])[:15]
            if cols and rows:
                data = [cols] + [[str(r.get(c, ""))[:60] for c in cols] for r in rows]
                t = RLTable(data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), HexColor("#001E44")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), white),
                    ("FONTNAME", (0, 0), (-1, 0), fn_bold),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("FONTNAME", (0, 1), (-1, -1), fn),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [white, HexColor("#F9FAFB")]),
                    ("GRID", (0, 0), (-1, -1), 0.25, HexColor("#e5e7eb")),
                ]))
                flow.append(t)

    takeaways = report_structure.get("key_takeaways") or []
    if len(qa_pairs) >= 3 and takeaways:
        flow.append(PageBreak())
        flow.append(Paragraph("Key Takeaways", h2))
        for t in takeaways:
            flow.append(Paragraph(f"• {_strip_emphasis(t)}", body))

    doc.build(flow)
    return buf.getvalue()
